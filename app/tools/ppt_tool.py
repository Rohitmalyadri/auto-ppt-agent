"""
MCP Tool: PowerPoint Operations
--------------------------------
All PPT manipulation functions are modular MCP-style tools.
Each function can be called independently or via the agent loop.
"""

import os
import logging
from pathlib import Path
from typing import List, Optional

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

logger = logging.getLogger(__name__)

# ─── Module-level presentation state ──────────────────────────────────────────
_presentation: Optional[Presentation] = None
_output_path: Optional[str] = None

# ─── Design Constants ──────────────────────────────────────────────────────────
THEME = {
    "bg_dark":      RGBColor(0x0D, 0x10, 0x17),   # near-black
    "bg_slide":     RGBColor(0x13, 0x18, 0x25),   # deep navy
    "accent":       RGBColor(0x5E, 0x81, 0xF4),   # electric blue
    "accent2":      RGBColor(0x9B, 0x59, 0xB6),   # violet
    "text_primary": RGBColor(0xF0, 0xF4, 0xFF),   # off-white
    "text_muted":   RGBColor(0xA0, 0xAB, 0xC4),   # muted blue-grey
    "bullet_dot":   RGBColor(0x5E, 0x81, 0xF4),   # accent blue
}

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


# ══════════════════════════════════════════════════════════════════════════════
# MCP Tool 1: create_presentation
# ══════════════════════════════════════════════════════════════════════════════
def create_presentation(filename: str) -> dict:
    """
    MCP Tool — Initialize a new Presentation object and set output path.

    Args:
        filename: Output filename (with or without .pptx extension).

    Returns:
        dict with status and path.
    """
    global _presentation, _output_path

    try:
        if not filename.endswith(".pptx"):
            filename += ".pptx"

        output_dir = Path(__file__).parent.parent.parent / "outputs"
        output_dir.mkdir(parents=True, exist_ok=True)
        _output_path = str(output_dir / filename)

        _presentation = Presentation()
        _presentation.slide_width = SLIDE_W
        _presentation.slide_height = SLIDE_H

        logger.info(f"[PPT Tool] Presentation created → {_output_path}")
        return {"status": "ok", "path": _output_path, "message": f"Presentation initialized: {filename}"}

    except Exception as e:
        logger.error(f"[PPT Tool] create_presentation failed: {e}")
        return {"status": "error", "message": str(e)}


# ══════════════════════════════════════════════════════════════════════════════
# MCP Tool 2: add_slide
# ══════════════════════════════════════════════════════════════════════════════
def add_slide(title: str, bullets: List[str], slide_number: int = 0) -> dict:
    """
    MCP Tool — Add a content slide with title and bullet points.

    Args:
        title: Slide title text.
        bullets: List of 3–5 bullet point strings.
        slide_number: Slide index (for logging).

    Returns:
        dict with status.
    """
    global _presentation

    if _presentation is None:
        return {"status": "error", "message": "Presentation not initialized. Call create_presentation first."}

    try:
        slide_layout = _presentation.slide_layouts[6]  # blank layout
        slide = _presentation.slides.add_slide(slide_layout)

        # ── Background ──────────────────────────────────────────────────────
        _set_slide_background(slide, THEME["bg_slide"])

        # ── Accent bar (left edge) ───────────────────────────────────────────
        _add_accent_bar(slide)

        # ── Title ────────────────────────────────────────────────────────────
        _add_title(slide, title)

        # ── Bullets ──────────────────────────────────────────────────────────
        _add_bullets(slide, bullets)

        # ── Slide number chip ────────────────────────────────────────────────
        if slide_number > 0:
            _add_slide_number(slide, slide_number)

        logger.info(f"[PPT Tool] Text slide added: '{title}'")
        return {"status": "ok", "message": f"Slide added: '{title}'"}

    except Exception as e:
        logger.error(f"[PPT Tool] add_slide failed: {e}")
        return {"status": "error", "message": str(e)}


# ══════════════════════════════════════════════════════════════════════════════
# MCP Tool 3: add_image_slide
# ══════════════════════════════════════════════════════════════════════════════
def add_image_slide(title: str, image_path: str, caption: str = "") -> dict:
    """
    MCP Tool — Add a slide with an image (and optional caption).

    Args:
        title: Slide title.
        image_path: Absolute path to image file.
        caption: Optional caption text below the image.

    Returns:
        dict with status.
    """
    global _presentation

    if _presentation is None:
        return {"status": "error", "message": "Presentation not initialized."}

    if not Path(image_path).exists():
        logger.warning(f"[PPT Tool] Image not found: {image_path}. Adding text slide instead.")
        return add_slide(title, [caption or "Image placeholder — content not available."])

    try:
        slide_layout = _presentation.slide_layouts[6]
        slide = _presentation.slides.add_slide(slide_layout)

        _set_slide_background(slide, THEME["bg_dark"])
        _add_accent_bar(slide)
        _add_title(slide, title)

        # Image centered in lower portion
        img_left = Inches(1.2)
        img_top  = Inches(1.6)
        img_w    = Inches(10.9)
        img_h    = Inches(5.2)
        pic = slide.shapes.add_picture(image_path, img_left, img_top, img_w, img_h)

        # Caption below image (if provided)
        if caption:
            txBox = slide.shapes.add_textbox(
                Inches(1.2), Inches(6.85), Inches(10.9), Inches(0.5)
            )
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = caption
            p.alignment = PP_ALIGN.CENTER
            p.runs[0].font.size = Pt(11)
            p.runs[0].font.color.rgb = THEME["text_muted"]
            p.runs[0].font.italic = True

        logger.info(f"[PPT Tool] Image slide added: '{title}'")
        return {"status": "ok", "message": f"Image slide added: '{title}'"}

    except Exception as e:
        logger.error(f"[PPT Tool] add_image_slide failed: {e}")
        # Graceful degradation
        return add_slide(title, [caption or "Visual content placeholder."])


# ══════════════════════════════════════════════════════════════════════════════
# MCP Tool 4: add_title_slide
# ══════════════════════════════════════════════════════════════════════════════
def add_title_slide(title: str, subtitle: str = "") -> dict:
    """
    MCP Tool — Add a hero/cover slide with large title and subtitle.

    Args:
        title: Main presentation title.
        subtitle: Optional subtitle or tagline.

    Returns:
        dict with status.
    """
    global _presentation

    if _presentation is None:
        return {"status": "error", "message": "Presentation not initialized."}

    try:
        slide_layout = _presentation.slide_layouts[6]
        slide = _presentation.slides.add_slide(slide_layout)

        # Full dark background
        _set_slide_background(slide, THEME["bg_dark"])

        # Wide gradient accent bar at top
        from pptx.util import Inches as In
        bar = slide.shapes.add_shape(
            1,  # MSO_SHAPE_TYPE.RECTANGLE
            0, 0, SLIDE_W, Inches(0.07)
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = THEME["accent"]
        bar.line.fill.background()

        # Title — large, centered, bold
        txBox = slide.shapes.add_textbox(
            Inches(1.0), Inches(2.2), Inches(11.3), Inches(2.0)
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = title
        run.font.size = Pt(48)
        run.font.bold = True
        run.font.color.rgb = THEME["text_primary"]
        p.alignment = PP_ALIGN.CENTER

        # Accent underline
        line = slide.shapes.add_shape(
            1, Inches(4.5), Inches(4.4), Inches(4.33), Inches(0.04)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = THEME["accent"]
        line.line.fill.background()

        # Subtitle
        if subtitle:
            stBox = slide.shapes.add_textbox(
                Inches(1.0), Inches(4.6), Inches(11.3), Inches(1.0)
            )
            stf = stBox.text_frame
            sp = stf.paragraphs[0]
            srun = sp.add_run()
            srun.text = subtitle
            srun.font.size = Pt(20)
            srun.font.color.rgb = THEME["text_muted"]
            srun.font.italic = True
            sp.alignment = PP_ALIGN.CENTER

        # Bottom tag
        tagBox = slide.shapes.add_textbox(
            Inches(0.5), Inches(6.9), Inches(12.33), Inches(0.4)
        )
        tp = tagBox.text_frame.paragraphs[0]
        trun = tp.add_run()
        trun.text = "Generated by Auto-PPT Agent  •  Powered by MCP + Hugging Face"
        trun.font.size = Pt(10)
        trun.font.color.rgb = THEME["accent"]
        tp.alignment = PP_ALIGN.CENTER

        logger.info(f"[PPT Tool] Title slide added: '{title}'")
        return {"status": "ok", "message": f"Title slide added: '{title}'"}

    except Exception as e:
        logger.error(f"[PPT Tool] add_title_slide failed: {e}")
        return {"status": "error", "message": str(e)}


# ══════════════════════════════════════════════════════════════════════════════
# MCP Tool 5: save_presentation
# ══════════════════════════════════════════════════════════════════════════════
def save_presentation() -> dict:
    """
    MCP Tool — Save the current presentation to disk.

    Returns:
        dict with status and saved path.
    """
    global _presentation, _output_path

    if _presentation is None:
        return {"status": "error", "message": "No presentation to save."}
    if _output_path is None:
        return {"status": "error", "message": "Output path not set."}

    try:
        _presentation.save(_output_path)
        logger.info(f"[PPT Tool] Presentation saved → {_output_path}")
        return {"status": "ok", "path": _output_path, "message": f"Saved: {_output_path}"}

    except Exception as e:
        logger.error(f"[PPT Tool] save_presentation failed: {e}")
        return {"status": "error", "message": str(e)}


# ══════════════════════════════════════════════════════════════════════════════
# Helper / Private Functions
# ══════════════════════════════════════════════════════════════════════════════

def _set_slide_background(slide, color: RGBColor):
    """Fill slide background with a solid color."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_accent_bar(slide):
    """Add a thin vertical accent bar on the left edge."""
    bar = slide.shapes.add_shape(
        1, 0, Inches(0.9), Inches(0.06), Inches(6.2)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = THEME["accent"]
    bar.line.fill.background()


def _add_title(slide, title: str):
    """Add styled title text to a slide."""
    txBox = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.25), Inches(12.5), Inches(0.85)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(32)
    run.font.bold = True
    run.font.color.rgb = THEME["text_primary"]


def _add_bullets(slide, bullets: List[str]):
    """Add bullet points to slide."""
    txBox = slide.shapes.add_textbox(
        Inches(0.55), Inches(1.25), Inches(12.5), Inches(5.8)
    )
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        # Bullet prefix
        run_dot = p.add_run()
        run_dot.text = "▸  "
        run_dot.font.color.rgb = THEME["bullet_dot"]
        run_dot.font.size = Pt(16)
        run_dot.font.bold = True

        # Bullet text
        run = p.add_run()
        run.text = bullet
        run.font.size = Pt(17)
        run.font.color.rgb = THEME["text_primary"]

        # Spacing between bullets
        from pptx.util import Pt as SPt
        p.space_after = SPt(8)


def _add_slide_number(slide, number: int):
    """Add a small slide number chip in the bottom-right corner."""
    numBox = slide.shapes.add_textbox(
        Inches(12.5), Inches(7.1), Inches(0.65), Inches(0.3)
    )
    p = numBox.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = str(number)
    run.font.size = Pt(10)
    run.font.color.rgb = THEME["accent"]
    p.alignment = PP_ALIGN.RIGHT
